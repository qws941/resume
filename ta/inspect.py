#!/usr/bin/env python3
"""ta.pptx 오버플로우 + 데이터 점검 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import re

prs = Presentation("/home/jclee/dev/ta/output/ta.pptx")

print("=" * 90)
print("슬라이드 크기:", round(prs.slide_width / 914400, 2), "x", round(prs.slide_height / 914400, 2), "inches")
print("=" * 90)

for si, slide in enumerate(prs.slides, 1):
    print(f"\n{'#'*90}")
    print(f"# SLIDE {si}")
    print(f"{'#'*90}")
    for shape in slide.shapes:
        if not shape.has_table:
            # 텍스트박스도 체크
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    w_in = round(shape.width / 914400, 2)
                    h_in = round(shape.height / 914400, 2)
                    print(f"\n[텍스트박스] name={shape.name!r} size={w_in}x{h_in}in")
                    print(f"  text: {txt[:100]}")
            continue

        tbl = shape.table
        tbl_name = shape.name
        tbl_w = shape.width  # EMU
        tbl_h = shape.height
        print(f"\n{'='*80}")
        print(f"[표] {tbl_name} — {tbl.rows.__len__()}행 x {len(tbl.columns)}열")
        print(f"  표 크기: {round(tbl_w/914400,2)}x{round(tbl_h/914400,2)} inches")
        print(f"  위치: left={round(shape.left/914400,2)} top={round(shape.top/914400,2)}")

        # 컬럼 너비
        col_widths = []
        for ci, col in enumerate(tbl.columns):
            w_in = round(col.width / 914400, 2)
            col_widths.append(w_in)
        print(f"  컬럼 너비(in): {col_widths}")

        # 행 높이
        row_heights = []
        for ri, row in enumerate(tbl.rows):
            h_in = round(row.height / 914400, 2)
            row_heights.append(h_in)
        print(f"  행 높이(in): {row_heights}")

        # 셀별 텍스트 + 오버플로우 체크
        overflow_issues = []
        data_dump = []
        for ri, row in enumerate(tbl.rows):
            for ci, cell in enumerate(row.cells):
                txt = cell.text.strip()
                if not txt:
                    continue

                col_w_in = col_widths[ci] if ci < len(col_widths) else 0
                row_h_in = row_heights[ri] if ri < len(row_heights) else 0

                # 머지 체크 — span_width 계산
                # cell의 실제 merge 범위 고려
                # python-pptx의 tc 속성으로 gridSpan 확인
                tc = cell._tc
                grid_span = int(tc.get('gridSpan', '1'))
                row_span = int(tc.get('rowSpan', '1'))
                if grid_span > 1:
                    actual_w = sum(col_widths[ci:ci+grid_span])
                else:
                    actual_w = col_w_in

                # 폰트 크기 추출
                font_sizes = []
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)
                max_font = max(font_sizes) if font_sizes else 0

                # 텍스트 줄 수
                lines = txt.split('\n')
                num_lines = len(lines)
                max_line_len = max(len(l) for l in lines)

                # 간이 오버플로우 추정:
                # 한글: 약 font_size * 0.014 inch per char
                # 영문: 약 font_size * 0.008 inch per char
                # 셀 마진 좌우 약 0.1in
                if max_font > 0 and actual_w > 0:
                    usable_w = actual_w - 0.15  # 마진
                    # 한글 비율 추정
                    korean_chars = len(re.findall(r'[가-힣]', txt))
                    total_chars = len(txt.replace('\n',''))
                    kr_ratio = korean_chars / total_chars if total_chars > 0 else 0
                    avg_char_w = max_font * (0.014 * kr_ratio + 0.008 * (1 - kr_ratio))

                    for line in lines:
                        line_len = len(line)
                        est_line_w = line_len * avg_char_w
                        if est_line_w > usable_w * 1.05:  # 5% 여유
                            overflow_issues.append({
                                'row': ri, 'col': ci,
                                'text': line[:60],
                                'est_width': round(est_line_w, 2),
                                'cell_width': round(usable_w, 2),
                                'font_pt': max_font,
                                'overflow_pct': round((est_line_w/usable_w - 1)*100, 1)
                            })

                    # 높이 오버플로우 체크
                    line_h = max_font * 1.4 / 72  # line height in inches (1.4x font)
                    est_total_h = num_lines * line_h + 0.1  # margin
                    actual_h = row_h_in * row_span if row_span > 1 else row_h_in
                    if est_total_h > actual_h * 1.1:
                        overflow_issues.append({
                            'row': ri, 'col': ci,
                            'text': f"[높이] {num_lines}줄",
                            'est_width': round(est_total_h, 2),
                            'cell_width': round(actual_h, 2),
                            'font_pt': max_font,
                            'overflow_pct': round((est_total_h/actual_h - 1)*100, 1)
                        })

                # 데이터 덤프
                data_dump.append(f"  [{ri},{ci}] gs={grid_span} rs={row_span} font={max_font}pt w={actual_w}in | {txt[:80]}")

        print(f"\n  --- 셀 데이터 ---")
        for d in data_dump:
            print(d)

        if overflow_issues:
            print(f"\n  ⚠️  오버플로우 의심 ({len(overflow_issues)}건):")
            for ov in overflow_issues:
                print(f"    [{ov['row']},{ov['col']}] font={ov['font_pt']}pt "
                      f"텍스트={ov['est_width']}in > 셀={ov['cell_width']}in "
                      f"(+{ov['overflow_pct']}%) | {ov['text']}")
        else:
            print(f"\n  ✅ 오버플로우 없음")

print(f"\n\n{'='*90}")
print("점검 완료")
