#!/usr/bin/env python3
"""ta.pptx 비주얼 개선 스크립트"""
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn, nsdecls
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
import copy, os

SRC = '/home/jclee/dev/ta/ta.pptx'
DST = '/home/jclee/dev/ta/output/ta.pptx'

# === Color Palette (Professional Navy) ===
NAVY      = RGBColor(0x1B, 0x36, 0x5F)   # 헤더 배경
NAVY_LIGHT= RGBColor(0x2C, 0x4A, 0x7C)   # 서브 헤더
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_BG   = RGBColor(0xF2, 0xF4, 0xF7)   # 짝수 행 배경
GRAY_LINE = RGBColor(0xD0, 0xD5, 0xDD)   # 테이블 선
ACCENT    = RGBColor(0x29, 0x70, 0xFF)    # 강조색


def set_cell_fill(cell, color):
    """셀 배경색 설정"""
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = tcPr.find(qn('a:solidFill'))
    if solidFill is not None:
        tcPr.remove(solidFill)
    # Remove noFill if present
    noFill = tcPr.find(qn('a:noFill'))
    if noFill is not None:
        tcPr.remove(noFill)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    srgb = etree.SubElement(sf, qn('a:srgbClr'))
    srgb.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')


def set_cell_border(cell, color=GRAY_LINE, width=Pt(0.5)):
    """셀 테두리 설정"""
    tcPr = cell._tc.get_or_add_tcPr()
    # Remove existing borders
    for tag in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        old = tcPr.find(qn(tag))
        if old is not None:
            tcPr.remove(old)

    for side in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = etree.SubElement(tcPr, qn(f'a:{side}'))
        ln.set('w', str(int(width)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        srgb = etree.SubElement(sf, qn('a:srgbClr'))
        srgb.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')


def style_cell_text(cell, font_name='KoPub돋움체 Medium', font_size=Pt(8),
                    bold=False, color=BLACK, alignment=PP_ALIGN.CENTER):
    """셀 텍스트 스타일링"""
    for para in cell.text_frame.paragraphs:
        para.alignment = alignment
        for run in para.runs:
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_margins(cell, top=Pt(2), bottom=Pt(2), left=Pt(4), right=Pt(4)):
    """셀 여백 설정"""
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.set('marT', str(int(top)))
    tcPr.set('marB', str(int(bottom)))
    tcPr.set('marL', str(int(left)))
    tcPr.set('marR', str(int(right)))


def style_header_row(table, row_idx, bg_color=NAVY, text_color=WHITE,
                     font_size=Pt(9), font_name='KoPub돋움체 Bold'):
    """헤더 행 스타일링"""
    for cell in table.rows[row_idx].cells:
        set_cell_fill(cell, bg_color)
        set_cell_border(cell, color=bg_color, width=Pt(0.5))
        style_cell_text(cell, font_name=font_name, font_size=font_size,
                       bold=True, color=text_color)
        set_cell_margins(cell)


def style_data_row(table, row_idx, alt=False):
    """데이터 행 스타일링 (alt=True이면 짝수행 배경)"""
    bg = GRAY_BG if alt else WHITE
    for cell in table.rows[row_idx].cells:
        set_cell_fill(cell, bg)
        set_cell_border(cell, color=GRAY_LINE, width=Pt(0.5))
        style_cell_text(cell, font_name='KoPub돋움체 Medium', font_size=Pt(8),
                       bold=False, color=BLACK)
        set_cell_margins(cell)


def is_row_empty(table, row_idx):
    """행이 비었는지 확인"""
    return all(cell.text.strip() == '' for cell in table.rows[row_idx].cells)


def process_profile_table(table):
    """표 149 (인적사항 5r x 9c) 스타일링"""
    for ri in range(len(table.rows)):
        for ci, cell in enumerate(table.rows[ri].cells):
            set_cell_border(cell, color=GRAY_LINE, width=Pt(0.5))
            set_cell_margins(cell)
            txt = cell.text.strip()
            # 라벨 셀 (성명, 소속회사, 부서, 최종학력, 기술등급, 생년월일, 입사일 등)
            labels = ['성명', '소속회사', '부 서', '최종학력', '기술등급',
                      '생년월일', '입사일', '직위', '졸업년월', '참여율',
                      '본 사업', 'IT 근무경력', '전문기술']
            is_label = any(txt.startswith(l) for l in labels)
            if is_label:
                set_cell_fill(cell, NAVY)
                style_cell_text(cell, font_name='KoPub돋움체 Bold', font_size=Pt(8),
                               bold=True, color=WHITE)
            else:
                set_cell_fill(cell, WHITE)
                style_cell_text(cell, font_name='KoPub돋움체 Medium', font_size=Pt(9),
                               bold=False, color=BLACK)


def process_company_table(table):
    """표 150 (회사 근무이력 8r x 3c) 스타일링"""
    # 헤더
    style_header_row(table, 0, bg_color=NAVY, font_size=Pt(8))
    # 데이터 행
    for ri in range(1, len(table.rows)):
        if not is_row_empty(table, ri):
            style_data_row(table, ri, alt=(ri % 2 == 0))
        else:
            # 빈 행은 배경만
            for cell in table.rows[ri].cells:
                set_cell_fill(cell, WHITE)
                set_cell_border(cell, color=GRAY_LINE, width=Pt(0.5))


def process_cert_table(table):
    """표 165 (자격증 5r x 2c) 스타일링"""
    style_header_row(table, 0, bg_color=NAVY, font_size=Pt(9))
    for ri in range(1, len(table.rows)):
        if not is_row_empty(table, ri):
            style_data_row(table, ri, alt=(ri % 2 == 0))


def process_project_table(table):
    """표 148 / 표 1 (프로젝트 이력) 스타일링"""
    # Row 0: 개인이력 (merged header)
    style_header_row(table, 0, bg_color=NAVY, font_size=Pt(10))
    # Row 1: 칼럼 헤더
    style_header_row(table, 1, bg_color=NAVY_LIGHT, font_size=Pt(8))
    # Data rows
    for ri in range(2, len(table.rows)):
        if not is_row_empty(table, ri):
            style_data_row(table, ri, alt=(ri % 2 == 1))
            # 첫 번째 칼럼 (수행기간) 약간 강조
            cell0 = table.rows[ri].cells[0]
            style_cell_text(cell0, font_name='KoPub돋움체 Medium', font_size=Pt(8),
                           bold=False, color=RGBColor(0x1B, 0x36, 0x5F))
        else:
            for cell in table.rows[ri].cells:
                set_cell_fill(cell, WHITE)
                set_cell_border(cell, color=GRAY_LINE, width=Pt(0.5))


def main():
    prs = Presentation(SRC)

    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            name = shape.name
            rows = len(table.rows)
            cols = len(table.columns)

            print(f"  Processing: Slide {si+1}, {name} ({rows}r x {cols}c)")

            # 분류 by 구조
            if name == '표 149' and cols == 9:
                process_profile_table(table)
            elif name == '표 150' and cols == 3:
                process_company_table(table)
            elif name == '표 165' and cols == 2:
                process_cert_table(table)
            elif name in ('표 148', '표 1') and cols == 6:
                process_project_table(table)
            else:
                print(f"    (skipped unknown table)")

    prs.save(DST)
    print(f"\n✅ 저장 완료: {DST}")
    print(f"   크기: {os.path.getsize(DST):,} bytes")


if __name__ == '__main__':
    main()
