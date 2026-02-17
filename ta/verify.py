#!/usr/bin/env python3
"""PPTX 검증 스크립트 — 파일 간 일관성 + 경력 데이터 정합성."""

import re
import sys
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation

# ──────────────────────────────────────────────
# 1. 데이터 추출
# ──────────────────────────────────────────────

PPTX_DIR = Path(__file__).parent
FILES = sorted(PPTX_DIR.glob("*.pptx"))

# 경력 타임라인 (ground truth)
CAREER_TIMELINE = [
    ("한국항공우주산업(KAI)", "2018.10", "2019.10"),
    ("메타넷엠플랫폼", "2019.12", "2021.08"),
    ("국민대학교", "2021.09", "2022.04"),
    ("펀엔씨", "2022.05", "2022.07"),
    ("콴텍투자일임", "2022.08", "2024.03"),
    ("가온누리정보시스템", "2024.03", "2025.02"),
    ("아이티센 CTS", "2025.03", "현재"),
]


def parse_ym(s: str) -> date:
    """'2018.10' 또는 '현재' → date."""
    if s in ("현재", "present", ""):
        return date.today().replace(day=1)
    parts = s.strip().split(".")
    return date(int(parts[0]), int(parts[1]), 1)


def months_between(start: date, end: date) -> int:
    return (end.year - start.year) * 12 + (end.month - start.month)


def calc_career_months() -> int:
    """CAREER_TIMELINE 기반 총 IT 경력 월수 (중복 기간 제거)."""
    intervals: list[tuple[date, date]] = []
    for _, s, e in CAREER_TIMELINE:
        intervals.append((parse_ym(s), parse_ym(e)))
    # merge overlapping
    intervals.sort()
    merged: list[tuple[date, date]] = [intervals[0]]
    for start, end in intervals[1:]:
        if start <= merged[-1][1]:
            merged[-1] = (merged[-1][0], max(merged[-1][1], end))
        else:
            merged.append((start, end))
    return sum(months_between(s, e) + 1 for s, e in merged)


def extract_tables(pptx_path: Path) -> list[dict[str, Any]]:
    """pptx 내 모든 테이블 추출 → [{name, rows, cols, data}]."""
    prs = Presentation(str(pptx_path))
    tables = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            data = []
            for row in tbl.rows:
                data.append([cell.text.strip() for cell in row.cells])
            tables.append({
                "name": shape.name,
                "slide": slide_idx + 1,
                "rows": len(tbl.rows),
                "cols": len(tbl.columns),
                "data": data,
            })
    return tables


def extract_profile(tables: list[dict]) -> dict[str, str]:
    """인적사항 테이블에서 key-value 추출."""
    profile: dict[str, str] = {}
    # 인적사항은 보통 첫 번째 테이블, 9열 구조
    for tbl in tables:
        flat = []
        for row in tbl["data"]:
            flat.extend(row)
        for i, cell in enumerate(flat):
            if cell in ("성명", "이름", "성 명"):
                if i + 1 < len(flat):
                    profile["성명"] = flat[i + 1]
            elif "생년월일" in cell:
                if i + 1 < len(flat):
                    profile["생년월일"] = flat[i + 1]
            elif "소속회사" in cell or "소속 회사" in cell:
                if i + 1 < len(flat):
                    profile["소속회사"] = flat[i + 1]
            elif "IT" in cell and ("경력" in cell or "근무" in cell):
                if i + 1 < len(flat):
                    profile["IT경력"] = flat[i + 1]
            elif "학력" in cell:
                if i + 1 < len(flat):
                    profile["학력"] = flat[i + 1]
    return profile


def extract_companies(tables: list[dict]) -> list[dict[str, str]]:
    """회사 근무이력 테이블 추출."""
    companies = []
    for tbl in tables:
        # 3열 구조: 회사명 | 기간 | 역할 등
        if tbl["cols"] == 3:
            for row in tbl["data"]:
                # 헤더 행 스킵
                if any(h in row[0] for h in ("근무처", "회사명", "소속", "기간")):
                    continue
                if row[0] and not all(c == "" for c in row):
                    companies.append({
                        "회사": row[0],
                        "기간": row[1] if len(row) > 1 else "",
                        "비고": row[2] if len(row) > 2 else "",
                    })
    return companies


def extract_career_duration_text(tables: list[dict]) -> str | None:
    """IT 근무경력 텍스트 값 추출."""
    for tbl in tables:
        for row in tbl["data"]:
            for i, cell in enumerate(row):
                if "IT" in cell and ("경력" in cell or "근무" in cell):
                    if i + 1 < len(row):
                        return row[i + 1]
    return None


def parse_duration_text(text: str) -> int | None:
    """'7년 4개월' → 88 (months)."""
    if not text:
        return None
    m = re.search(r"(\d+)\s*년\s*(\d+)\s*개월", text)
    if m:
        return int(m.group(1)) * 12 + int(m.group(2))
    m = re.search(r"(\d+)\s*년", text)
    if m:
        return int(m.group(1)) * 12
    return None


# ──────────────────────────────────────────────
# 2. 검증 실행
# ──────────────────────────────────────────────

def run_verification():
    print("=" * 72)
    print("  PPTX 검증 리포트")
    print(f"  대상: {len(FILES)}개 파일  |  기준일: {date.today()}")
    print("=" * 72)

    all_data: dict[str, dict[str, Any]] = {}

    # ── Phase 1: 전체 추출 ──
    print("\n[Phase 1] 데이터 추출")
    print("-" * 40)
    for f in FILES:
        fname = f.name
        try:
            tables = extract_tables(f)
            profile = extract_profile(tables)
            companies = extract_companies(tables)
            duration_text = extract_career_duration_text(tables)
            all_data[fname] = {
                "tables": tables,
                "profile": profile,
                "companies": companies,
                "duration_text": duration_text,
                "table_count": len(tables),
            }
            print(f"  ✓ {fname}: {len(tables)} tables, "
                  f"profile={bool(profile)}, companies={len(companies)}")
        except Exception as e:
            all_data[fname] = {"error": str(e)}
            print(f"  ✗ {fname}: ERROR — {e}")

    # ── Phase 2: 파일 간 일관성 비교 ──
    print(f"\n{'=' * 72}")
    print("[Phase 2] 파일 간 일관성 비교")
    print("-" * 40)

    valid_files = {k: v for k, v in all_data.items() if "error" not in v}
    issues: list[str] = []

    # 2a. 성명 비교
    names = {k: v["profile"].get("성명", "N/A") for k, v in valid_files.items()}
    unique_names = set(names.values()) - {"N/A", ""}
    if len(unique_names) > 1:
        issues.append(f"성명 불일치: {names}")
    print(f"  성명: {unique_names or 'N/A'} — {'✓ 일치' if len(unique_names) <= 1 else '✗ 불일치'}")

    # 2b. 소속회사 비교
    orgs = {k: v["profile"].get("소속회사", "N/A") for k, v in valid_files.items()}
    unique_orgs = set(orgs.values()) - {"N/A", ""}
    if len(unique_orgs) > 1:
        issues.append(f"소속회사 불일치: {orgs}")
        print(f"  소속회사: ✗ 불일치")
        for fname, org in orgs.items():
            if org and org != "N/A":
                print(f"    - {fname}: {org}")
    else:
        print(f"  소속회사: {unique_orgs or 'N/A'} — ✓ 일치")

    # 2c. IT경력 비교
    durations = {k: v.get("duration_text", "N/A") for k, v in valid_files.items()}
    unique_durations = set(d for d in durations.values() if d and d != "N/A")
    if len(unique_durations) > 1:
        issues.append(f"IT경력 불일치: {durations}")
        print(f"  IT경력: ✗ 불일치")
        for fname, dur in durations.items():
            if dur and dur != "N/A":
                print(f"    - {fname}: {dur}")
    else:
        print(f"  IT경력: {unique_durations or 'N/A'} — {'✓ 일치' if len(unique_durations) <= 1 else '?'}")

    # 2d. 회사이력 수 비교
    company_counts = {k: len(v["companies"]) for k, v in valid_files.items()}
    if len(set(company_counts.values())) > 1:
        issues.append(f"회사이력 수 불일치: {company_counts}")
        print(f"  회사이력: ✗ 파일별 수 불일치")
        for fname, cnt in company_counts.items():
            print(f"    - {fname}: {cnt}개사")
    else:
        v = list(company_counts.values())[0] if company_counts else 0
        print(f"  회사이력: {v}개사 — ✓ 일치")

    # 2e. 생년월일 비교
    births = {k: v["profile"].get("생년월일", "N/A") for k, v in valid_files.items()}
    unique_births = set(b for b in births.values() if b and b != "N/A")
    if len(unique_births) > 1:
        issues.append(f"생년월일 불일치: {births}")
    print(f"  생년월일: {unique_births or 'N/A'} — {'✓ 일치' if len(unique_births) <= 1 else '✗ 불일치'}")

    # ── Phase 3: 경력 정합성 검증 ──
    print(f"\n{'=' * 72}")
    print("[Phase 3] 경력 데이터 정합성 검증")
    print("-" * 40)

    # 3a. Ground truth 경력 계산
    actual_months = calc_career_months()
    actual_y, actual_m = divmod(actual_months, 12)
    print(f"  Ground truth 경력: {actual_y}년 {actual_m}개월 ({actual_months}개월)")
    print(f"  타임라인: {CAREER_TIMELINE[0][1]} ~ 현재")
    print()

    # 3b. 각 파일의 기재 경력 vs 실제
    for fname, data in valid_files.items():
        dt = data.get("duration_text")
        if not dt:
            continue
        file_months = parse_duration_text(dt)
        if file_months is None:
            print(f"  {fname}: '{dt}' — 파싱 불가")
            continue
        diff = file_months - actual_months
        status = "✓ 정확" if abs(diff) <= 1 else f"✗ {'+' if diff > 0 else ''}{diff}개월 차이"
        print(f"  {fname}: 기재 '{dt}' ({file_months}개월) — {status}")
        if abs(diff) > 1:
            issues.append(f"{fname}: IT경력 {diff:+d}개월 오차 (기재={dt}, 실제={actual_y}년 {actual_m}개월)")

    # 3c. 회사이력 누락 검증
    print()
    print("  회사이력 누락 검증:")
    gt_companies = {c[0] for c in CAREER_TIMELINE}
    for fname, data in valid_files.items():
        file_companies = {c["회사"] for c in data["companies"]}
        if not file_companies:
            print(f"    {fname}: 회사이력 테이블 없음 (또는 빈 테이블)")
            continue
        missing = gt_companies - file_companies
        # fuzzy match — 부분 문자열 매치 시 제외
        real_missing = set()
        for mc in missing:
            if not any(mc[:3] in fc or fc[:3] in mc for fc in file_companies):
                real_missing.add(mc)
        if real_missing:
            issues.append(f"{fname}: 회사이력 누락 — {real_missing}")
            print(f"    {fname}: ✗ 누락 {len(real_missing)}개 — {', '.join(sorted(real_missing))}")
        else:
            print(f"    {fname}: ✓ 전체 포함 ({len(file_companies)}개사)")

    # 3d. 경력 공백 기간 분석
    print()
    print("  경력 공백 기간:")
    for i in range(1, len(CAREER_TIMELINE)):
        prev_end = parse_ym(CAREER_TIMELINE[i - 1][2])
        curr_start = parse_ym(CAREER_TIMELINE[i][1])
        gap = months_between(prev_end, curr_start)
        if gap > 1:
            print(f"    ⚠ {CAREER_TIMELINE[i-1][0]} → {CAREER_TIMELINE[i][0]}: "
                  f"{gap}개월 공백 ({CAREER_TIMELINE[i-1][2]} ~ {CAREER_TIMELINE[i][1]})")
        elif gap == 0:
            print(f"    ↔ {CAREER_TIMELINE[i-1][0]} → {CAREER_TIMELINE[i][0]}: 중복/연속")

    # ── 요약 ──
    print(f"\n{'=' * 72}")
    print("[요약]")
    print("-" * 40)
    if issues:
        print(f"  ✗ {len(issues)}개 이슈 발견:")
        for i, iss in enumerate(issues, 1):
            print(f"    {i}. {iss}")
    else:
        print("  ✓ 모든 검증 통과 — 이슈 없음")
    print(f"{'=' * 72}")

    return len(issues)


if __name__ == "__main__":
    sys.exit(run_verification())
