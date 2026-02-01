#!/usr/bin/env python3
"""
PPTX Generation Engine - Common logic extracted from generate_ta/shinhan_pptx.py

Usage:
    from pptx_engine import generate
    from pptx_templates import TEMPLATES
    generate(TEMPLATES["ta"], source_data)
"""

import datetime
import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable

try:
    from pptx import Presentation
except ImportError:
    import sys
    print("python-pptx required: pip install python-pptx")
    sys.exit(1)

from pptx.util import Pt
from pptx_utils import apply_korean_font_to_table

logging.basicConfig(level=logging.INFO, format="%(message)s")
logger = logging.getLogger(__name__)


@dataclass
class TemplateSpec:
    """Template specification for PPTX generation."""
    name: str
    source_path: Path
    template_path: Path
    output_path: Path
    # Maps (rows, cols) -> handler function
    table_handlers: dict = field(default_factory=dict)
    # Optional: title font size (default 24pt)
    title_size_pt: int = 24


def normalize_slide_title(slide, size_pt: int = 24) -> bool:
    """
    Normalize slide title to specified font size.
    Returns True if title was found and normalized.
    
    Strategy:
    1. Use slide.shapes.title if available
    2. Fallback: find topmost text shape within 100pt from top
    """
    # Try official title shape first
    if slide.shapes.title:
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(size_pt)
        return True
    
    # Fallback: find topmost text shape
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    if text_shapes:
        topmost = min(text_shapes, key=lambda s: s.top)
        if topmost.top < Pt(100):
            for p in topmost.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(size_pt)
            return True
    
    return False


def get_current_age(birth_date_str: str) -> int:
    """Calculate current age from birth date string (YYYY.MM.DD or YYYY-MM-DD)."""
    birth_year = int(birth_date_str.split(".")[0].split("-")[0])
    return datetime.date.today().year - birth_year


def generate(spec: TemplateSpec, source: dict = None) -> Path:
    """
    Generate PPTX from template using source data.
    
    Args:
        spec: TemplateSpec with paths and table handlers
        source: Optional source data dict. If None, loads from spec.source_path
    
    Returns:
        Path to generated output file
    """
    # Load source if not provided
    if source is None:
        logger.info(f"📄 Loading: {spec.source_path}")
        with open(spec.source_path, "r", encoding="utf-8") as f:
            source = json.load(f)
    
    logger.info(f"📝 Filling: {spec.template_path}")
    prs = Presentation(str(spec.template_path))
    
    unhandled_tables = []
    
    for slide_idx, slide in enumerate(prs.slides):
        # Normalize title
        normalize_slide_title(slide, spec.title_size_pt)
        
        # Process tables
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            
            tbl = shape.table
            rows, cols = len(tbl.rows), len(tbl.columns)
            key = (rows, cols)
            
            handler = spec.table_handlers.get(key)
            if handler:
                handler(tbl, source)
                apply_korean_font_to_table(tbl)
            else:
                unhandled_tables.append(f"Slide {slide_idx + 1}: {rows}x{cols}")
    
    # Warn about unhandled tables
    if unhandled_tables:
        logger.warning(f"⚠️  Unhandled tables: {', '.join(unhandled_tables)}")
    
    # Save output
    spec.output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(spec.output_path))
    
    logger.info(f"✅ Generated: {spec.output_path}")
    if "personal" in source:
        logger.info(f"   - Name: {source['personal'].get('name', 'N/A')}")
    if "summary" in source:
        logger.info(f"   - Experience: {source['summary'].get('totalExperience', 'N/A')}")
    
    return spec.output_path


def generate_from_cli(template_name: str, source_path: Path = None, output_path: Path = None) -> Path:
    """
    CLI wrapper for generate().
    
    Args:
        template_name: Template name ("ta" or "shinhan")
        source_path: Optional override for source data path
        output_path: Optional override for output path
    
    Returns:
        Path to generated output file
    """
    # Import here to avoid circular dependency
    from pptx_templates import TEMPLATES
    
    if template_name not in TEMPLATES:
        raise ValueError(f"Unknown template: {template_name}. Available: {list(TEMPLATES.keys())}")
    
    spec = TEMPLATES[template_name]
    
    # Override paths if provided
    if source_path:
        spec = TemplateSpec(
            name=spec.name,
            source_path=source_path,
            template_path=spec.template_path,
            output_path=output_path or spec.output_path,
            table_handlers=spec.table_handlers,
            title_size_pt=spec.title_size_pt,
        )
    elif output_path:
        spec = TemplateSpec(
            name=spec.name,
            source_path=spec.source_path,
            template_path=spec.template_path,
            output_path=output_path,
            table_handlers=spec.table_handlers,
            title_size_pt=spec.title_size_pt,
        )
    
    return generate(spec)
