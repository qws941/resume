#!/usr/bin/env python3
"""
PPTX Debugging Tool - Unified CLI for PPTX inspection and maintenance

Usage:
    python pptx_tool.py verify [file]        # Verify font sizes (default: ta_filled.pptx)
    python pptx_tool.py dump [file]          # Dump all text content
    python pptx_tool.py titles [file]        # Find title shapes and their sizes
    python pptx_tool.py fix-typos [file]     # Fix common typos
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("python-pptx required: pip install python-pptx")
    sys.exit(1)

ROOT = Path(__file__).parent.parent.parent
DEFAULT_FILE = ROOT / "resumes/generated/ta_filled.pptx"

TYPO_REPLACEMENTS = {
    "Complience": "Compliance",
    "Architecure": "Architecture",
}


def cmd_verify(args):
    prs = Presentation(str(args.file))
    expected_body_min = args.min_body or 13.5
    
    print(f"\n--- Verifying: {args.file} (min body: {expected_body_min}pt) ---")
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"\nSlide {i}:")
        
        title_size = _get_title_size(slide)
        if title_size:
            status = "PASS" if abs(title_size - 24) < 0.5 else "FAIL"
            print(f"  Title: {title_size}pt [{status}]")
        
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            cell_00_size = _get_cell_font_size(tbl.cell(0, 0))
            if cell_00_size:
                status = "PASS" if abs(cell_00_size - 24) < 0.5 else "FAIL"
                print(f"  Table (0,0): {cell_00_size}pt [{status}]")
            
            body_sizes = []
            for r_idx, row in enumerate(tbl.rows):
                for c_idx, cell in enumerate(row.cells):
                    if r_idx == 0 and c_idx == 0:
                        continue
                    size = _get_cell_font_size(cell)
                    if size:
                        body_sizes.append(size)
            
            if body_sizes:
                min_body = min(body_sizes)
                max_body = max(body_sizes)
                avg_body = sum(body_sizes) / len(body_sizes)
                print(f"  Table body: avg={avg_body:.1f}pt min={min_body}pt max={max_body}pt")
                if min_body < expected_body_min:
                    print(f"  ⚠️  WARNING: Found text smaller than {expected_body_min}pt")


def cmd_dump(args):
    prs = Presentation(str(args.file))
    print(f"\n--- All Text in: {args.file} ---")
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"\nSlide {i}:")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text[:80].replace("\n", " ")
                print(f"  [{shape.name}]: \"{text}\"")


def cmd_titles(args):
    prs = Presentation(str(args.file))
    print(f"\n--- Titles in: {args.file} ---")
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"\nSlide {i}:")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if not text or len(text) > 50:
                continue
            sizes = []
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        sizes.append(r.font.size.pt)
            print(f"  [{shape.name}]: \"{text}\" | sizes={sizes}")


def cmd_fix_typos(args):
    prs = Presentation(str(args.file))
    print(f"\n--- Fixing typos in: {args.file} ---")
    
    total_fixes = 0
    
    for i, slide in enumerate(prs.slides, 1):
        slide_fixes = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        for wrong, correct in TYPO_REPLACEMENTS.items():
                            if wrong in r.text:
                                r.text = r.text.replace(wrong, correct)
                                slide_fixes += 1
                                print(f"  Slide {i}: '{wrong}' → '{correct}'")
            
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                for wrong, correct in TYPO_REPLACEMENTS.items():
                                    if wrong in r.text:
                                        r.text = r.text.replace(wrong, correct)
                                        slide_fixes += 1
                                        print(f"  Slide {i} table: '{wrong}' → '{correct}'")
        
        total_fixes += slide_fixes
    
    print(f"\n✅ Total fixes: {total_fixes}")
    
    if total_fixes > 0:
        output = args.output or args.file
        prs.save(str(output))
        print(f"💾 Saved: {output}")


def _get_title_size(slide):
    if slide.shapes.title:
        for p in slide.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    return r.font.size.pt
    
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    if text_shapes:
        topmost = min(text_shapes, key=lambda s: s.top)
        if topmost.top < Pt(100):
            for p in topmost.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        return r.font.size.pt
    return None


def _get_cell_font_size(cell):
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                return r.font.size.pt
    return None


def main():
    parser = argparse.ArgumentParser(description="PPTX debugging and maintenance tool")
    subparsers = parser.add_subparsers(dest="command", required=True)
    
    p_verify = subparsers.add_parser("verify", help="Verify font sizes")
    p_verify.add_argument("file", nargs="?", type=Path, default=DEFAULT_FILE)
    p_verify.add_argument("--min-body", type=float, help="Minimum body font size (default: 13.5)")
    p_verify.set_defaults(func=cmd_verify)
    
    p_dump = subparsers.add_parser("dump", help="Dump all text content")
    p_dump.add_argument("file", nargs="?", type=Path, default=DEFAULT_FILE)
    p_dump.set_defaults(func=cmd_dump)
    
    p_titles = subparsers.add_parser("titles", help="Find title shapes")
    p_titles.add_argument("file", nargs="?", type=Path, default=DEFAULT_FILE)
    p_titles.set_defaults(func=cmd_titles)
    
    p_fix = subparsers.add_parser("fix-typos", help="Fix common typos")
    p_fix.add_argument("file", nargs="?", type=Path, default=DEFAULT_FILE)
    p_fix.add_argument("--output", "-o", type=Path, help="Output file (default: overwrite)")
    p_fix.set_defaults(func=cmd_fix_typos)
    
    args = parser.parse_args()
    
    if not args.file.exists():
        print(f"❌ File not found: {args.file}")
        return 1
    
    args.func(args)
    return 0


if __name__ == "__main__":
    sys.exit(main())
