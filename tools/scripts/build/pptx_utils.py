#!/usr/bin/env python3

from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

LIMITS = {
    "company": 20,
    "description": 70,  # 50→70: 담당업무 잘림 해결
    "technologies": 45,  # 35→45: 기술스택 잘림 해결
    "expertise": 40,    # 30→40: 전문분야 여유 확보
    "project_name": 25,
    "default": 50,
}

DEFAULT_KOREAN_FONT = "NanumGothic"


def truncate(text, field=None, max_len=None):
    limit = max_len or (LIMITS.get(field, LIMITS["default"]) if field else LIMITS["default"])
    text = str(text)
    return text[:limit-2] + ".." if len(text) > limit else text


def join_truncate(items, field=None, max_len=None, max_items=None):
    if max_items:
        items = items[:max_items]
    joined = ", ".join(str(i) for i in items)
    return truncate(joined, field, max_len)


def resize_table_text(table, header_rows=0, header_cols=0, header_size_pt=24, body_size_pt=18):
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            is_header = r_idx < header_rows or c_idx < header_cols
            size_pt = header_size_pt if is_header else body_size_pt
            
            if not cell.text_frame:
                continue
                
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(size_pt)


def set_cell_text(cell, text, field=None, max_len=None, font_size_pt=None):
    tf = cell.text_frame
    if not (tf.paragraphs and tf.paragraphs[0].runs):
        cell.text = str(text)
        if font_size_pt:
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size_pt)
        return

    first_para = tf.paragraphs[0]
    first_run = first_para.runs[0]

    font_name = first_run.font.name
    font_size = first_run.font.size
    font_bold = first_run.font.bold
    font_italic = first_run.font.italic
    font_color = None
    if first_run.font.color and first_run.font.color.type is not None:
        try:
            font_color = first_run.font.color.rgb
        except AttributeError:
            pass

    para_alignment = first_para.alignment

    for p in tf.paragraphs:
        p.clear()

    new_run = tf.paragraphs[0].add_run()
    final_text = truncate(text, field, max_len) if field or max_len else str(text)
    new_run.text = final_text

    if font_name:
        new_run.font.name = font_name
    
    if font_size_pt:
        new_run.font.size = Pt(font_size_pt)
    elif font_size:
        new_run.font.size = font_size
        
    if font_bold is not None:
        new_run.font.bold = font_bold
    if font_italic is not None:
        new_run.font.italic = font_italic
    if font_color:
        new_run.font.color.rgb = font_color
    if para_alignment:
        tf.paragraphs[0].alignment = para_alignment
    
    effective_font = font_name or DEFAULT_KOREAN_FONT
    set_east_asian_font(new_run, effective_font)


def _remove_font_elements(rPr, tag):
    for elem in rPr.findall(qn(tag)):
        rPr.remove(elem)


def _add_font_element(rPr, tag, font_name):
    elem = etree.SubElement(rPr, qn(tag))
    elem.set('typeface', font_name)


def set_east_asian_font(run, font_name):
    """python-pptx's run.font.name only sets a:latin, not a:ea (East Asian)."""
    rPr = run._r.get_or_add_rPr()
    
    _remove_font_elements(rPr, 'a:ea')
    _add_font_element(rPr, 'a:ea', font_name)
    
    _remove_font_elements(rPr, 'a:latin')
    _add_font_element(rPr, 'a:latin', font_name)


def apply_korean_font_to_table(table, font_name=None):
    font = font_name or DEFAULT_KOREAN_FONT
    for row in table.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        set_east_asian_font(run, font)


def hide_empty_rows(table, start_row=0, check_col=0):
    empty_indices = []
    for r_idx in range(len(table.rows) - 1, start_row - 1, -1):
        if not table.cell(r_idx, check_col).text_frame.text.strip():
            empty_indices.append(r_idx)
    
    tr_elements = table._tbl.findall(qn('a:tr'))
    for r_idx in empty_indices:
        if r_idx < len(tr_elements):
            tr_elements[r_idx].getparent().remove(tr_elements[r_idx])
    
    return len(empty_indices)
