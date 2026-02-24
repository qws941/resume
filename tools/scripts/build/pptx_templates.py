#!/usr/bin/env python3
"""
PPTX Template Specifications - Table handlers for each template type

Each handler function receives (table, source) and fills the table cells.
"""

from pathlib import Path
from pptx.util import Pt

try:
    from .pptx_engine import TemplateSpec, get_current_age
    from .pptx_utils import (
        set_cell_text,
        join_truncate,
        truncate,
        resize_table_text,
        hide_empty_rows,
    )
except ImportError:
    import importlib
    import sys

    build_dir = Path(__file__).resolve().parent
    if str(build_dir) not in sys.path:
        sys.path.insert(0, str(build_dir))

    _pptx_engine = importlib.import_module("pptx_engine")
    _pptx_utils = importlib.import_module("pptx_utils")

    TemplateSpec = _pptx_engine.TemplateSpec
    get_current_age = _pptx_engine.get_current_age
    set_cell_text = _pptx_utils.set_cell_text
    join_truncate = _pptx_utils.join_truncate
    truncate = _pptx_utils.truncate
    resize_table_text = _pptx_utils.resize_table_text
    hide_empty_rows = _pptx_utils.hide_empty_rows

ROOT = Path(__file__).parent.parent.parent.parent


# =============================================================================
# TA Template Handlers (15x8 profile + 12x6 projects)
# =============================================================================


def handle_ta_profile(tbl, source):
    """TA 기본사항 (15x8): 인적사항 + 경력사항 + 자격증"""
    resize_table_text(
        tbl, header_rows=2, header_cols=1, header_size_pt=14, body_size_pt=14
    )

    # Force Title Cell (0,0) to 24pt
    if tbl.cell(0, 0).text_frame.text.strip():
        for p in tbl.cell(0, 0).text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(24)

    p = source["personal"]
    e = source["education"]
    s = source["summary"]

    # Row 2: 이름, 소속, 직위, 등급, 학력, 전공
    set_cell_text(tbl.cell(2, 0), p["name"])
    set_cell_text(tbl.cell(2, 1), source["current"]["company"].replace("(주)", ""))
    set_cell_text(tbl.cell(2, 2), source["current"]["position"])
    set_cell_text(tbl.cell(2, 4), s["grade"])
    set_cell_text(tbl.cell(2, 5), f"{e['school']}({e['status']})")
    set_cell_text(tbl.cell(2, 6), e["major"])

    # Row 3: 나이, 총경력
    age = get_current_age(p["birthDate"])
    set_cell_text(tbl.cell(3, 1), f"만   {age} 세")
    set_cell_text(tbl.cell(3, 4), s["totalExperience"])

    # Row 2 Col 7: 전문분야
    expertise_text = ", ".join(s.get("expertise", []))
    set_cell_text(tbl.cell(2, 7), expertise_text)

    # 경력사항 (Row 8~14)
    rows = len(tbl.rows)
    careers = source["careers"][:7]
    for i, career in enumerate(careers):
        row_idx = 8 + i
        if row_idx < rows:
            set_cell_text(tbl.cell(row_idx, 0), truncate(career["company"], "company"))
            set_cell_text(tbl.cell(row_idx, 1), career["period"])
            set_cell_text(
                tbl.cell(row_idx, 4), truncate(career["description"], "description")
            )

    # 자격증 (Row 8~14, Col 5=자격증명, Col 7=취득일)
    certs = source.get("certifications", [])[:7]
    for i, cert in enumerate(certs):
        row_idx = 8 + i
        if row_idx < rows:
            set_cell_text(tbl.cell(row_idx, 5), cert["name"])
            set_cell_text(tbl.cell(row_idx, 7), cert.get("date", ""))


def handle_ta_projects(tbl, source):
    """TA 프로젝트 참여이력 (12x6)"""
    resize_table_text(tbl, header_rows=2, header_size_pt=14, body_size_pt=14)

    # Force Title Cell (0,0) to 24pt
    if tbl.cell(0, 0).text_frame.text.strip():
        for p in tbl.cell(0, 0).text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(24)

    rows = len(tbl.rows)
    projects = source["projects"][:9]
    for i, proj in enumerate(projects):
        row_idx = 2 + i
        if row_idx < rows:
            set_cell_text(tbl.cell(row_idx, 0), truncate(proj["name"], "project_name"))
            set_cell_text(tbl.cell(row_idx, 1), proj["role"])
            set_cell_text(tbl.cell(row_idx, 2), proj.get("description", ""))
            set_cell_text(
                tbl.cell(row_idx, 3),
                join_truncate(proj["technologies"], "technologies", max_items=3),
            )
            set_cell_text(tbl.cell(row_idx, 4), proj["period"].replace(" ~ ", "~"))
            set_cell_text(tbl.cell(row_idx, 5), proj["client"])

    hide_empty_rows(tbl, start_row=2, check_col=0)


# =============================================================================
# Shinhan Template Handlers
# =============================================================================


def handle_shinhan_profile(tbl, source):
    """신한 기본사항 (5x9)"""
    resize_table_text(tbl, header_cols=1, header_size_pt=14, body_size_pt=12)

    p = source["personal"]
    c = source["current"]
    e = source["education"]
    s = source["summary"]

    set_cell_text(tbl.cell(0, 1), p["name"])
    set_cell_text(tbl.cell(0, 7), p["birthDate"])
    set_cell_text(tbl.cell(1, 1), c["company"])
    set_cell_text(tbl.cell(1, 7), c["startDate"])
    set_cell_text(tbl.cell(2, 1), c["department"])
    set_cell_text(tbl.cell(2, 3), c["position"])
    set_cell_text(tbl.cell(2, 5), str(c["participationMonths"]))
    set_cell_text(tbl.cell(2, 8), c["participationRate"])
    set_cell_text(tbl.cell(3, 1), f"{e['school']} {e['major']}")
    set_cell_text(tbl.cell(3, 5), e["status"])
    set_cell_text(tbl.cell(4, 1), s["grade"])
    set_cell_text(tbl.cell(4, 4), s["totalExperience"])
    set_cell_text(tbl.cell(4, 7), join_truncate(s["expertise"], "expertise"))


def handle_shinhan_careers(tbl, source):
    """신한 경력사항 (5x3)"""
    resize_table_text(tbl, header_rows=1, header_size_pt=14, body_size_pt=12)

    for i, career in enumerate(source["careers"][:4]):
        set_cell_text(tbl.cell(i + 1, 0), career["company"])
        set_cell_text(tbl.cell(i + 1, 1), career["period"].replace(" ~ ", "~"))
        set_cell_text(tbl.cell(i + 1, 2), career["project"])


def handle_shinhan_certs_small(tbl, source):
    """신한 자격증 (5x2)"""
    resize_table_text(tbl, header_rows=1, header_size_pt=14, body_size_pt=12)

    certs = source["certifications"][:4]
    for i, cert in enumerate(certs):
        set_cell_text(tbl.cell(i + 1, 0), cert["name"])
        set_cell_text(tbl.cell(i + 1, 1), cert["date"])


def handle_shinhan_certs_large(tbl, source):
    """신한 자격증 (7x2)"""
    resize_table_text(tbl, header_rows=1, header_size_pt=14, body_size_pt=12)

    rows = len(tbl.rows)
    certs = source["certifications"][:6]
    for i, cert in enumerate(certs):
        if i + 1 < rows:
            set_cell_text(tbl.cell(i + 1, 0), cert["name"])
            set_cell_text(tbl.cell(i + 1, 1), cert["date"])


def handle_shinhan_projects(tbl, source):
    """신한 프로젝트 (9x6)"""
    resize_table_text(tbl, header_rows=2, header_size_pt=14, body_size_pt=11)

    rows = len(tbl.rows)
    for i, proj in enumerate(source["projects"][:7]):
        if i + 2 < rows:
            set_cell_text(tbl.cell(i + 2, 0), proj["period"].replace(" ~ ", "~"))
            set_cell_text(tbl.cell(i + 2, 1), proj["name"])
            set_cell_text(tbl.cell(i + 2, 2), proj["client"])
            set_cell_text(
                tbl.cell(i + 2, 3), join_truncate(proj["technologies"], "technologies")
            )
            set_cell_text(tbl.cell(i + 2, 4), proj["os"])
            set_cell_text(tbl.cell(i + 2, 5), proj["role"])


def handle_shinhan_personal_projects(tbl, source):
    """신한 개인프로젝트 (11x6)"""
    resize_table_text(tbl, header_rows=2, header_size_pt=14, body_size_pt=11)

    rows = len(tbl.rows)
    extras = source.get("personalProjects", [])
    for i, proj in enumerate(extras[:5]):
        if i + 2 < rows:
            set_cell_text(tbl.cell(i + 2, 0), proj["period"].replace(" ~ ", "~"))
            set_cell_text(tbl.cell(i + 2, 1), proj["name"])
            set_cell_text(tbl.cell(i + 2, 2), "개인")
            set_cell_text(
                tbl.cell(i + 2, 3), join_truncate(proj["technologies"], "technologies")
            )
            set_cell_text(tbl.cell(i + 2, 4), "Docker")
            set_cell_text(tbl.cell(i + 2, 5), "설계/개발")


# =============================================================================
# Template Definitions
# =============================================================================

TEMPLATES = {
    "ta": TemplateSpec(
        name="TA형 이력서",
        source_path=ROOT / "typescript/data/resumes/master/resume_data.json",
        template_path=ROOT / "typescript/data/resumes/generated/ta.pptx",
        output_path=ROOT / "typescript/data/resumes/generated/ta_filled.pptx",
        table_handlers={
            (15, 8): handle_ta_profile,
            (12, 6): handle_ta_projects,
        },
    ),
    "shinhan": TemplateSpec(
        name="신한형 이력서",
        source_path=ROOT / "typescript/data/resumes/master/shinhan_resume_data.json",
        template_path=ROOT / "typescript/data/resumes/archive/기타/shinhan.pptx",
        output_path=ROOT / "typescript/data/resumes/generated/shinhan_filled.pptx",
        table_handlers={
            (5, 9): handle_shinhan_profile,
            (5, 3): handle_shinhan_careers,
            (5, 2): handle_shinhan_certs_small,
            (7, 2): handle_shinhan_certs_large,
            (9, 6): handle_shinhan_projects,
            (11, 6): handle_shinhan_personal_projects,
        },
    ),
}
